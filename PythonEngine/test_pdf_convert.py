"""Tự kiểm tra chuyển đổi PDF -> Word/PPTX. Chạy trực tiếp:
    python3 test_pdf_convert.py
"""
import os
import tempfile

import fitz
from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor

from pdf_convert import (
    convert_to_docx, convert_to_pptx, EMU_PER_POINT, _detect_table, _build_table_cells,
    _page_background_png, _table_bbox_rect, _add_table_shape, _table_crosses_section_divider,
)


def make_sample_pdf(path):
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 100), "Hello world, sample paragraph.", fontsize=14,
                      color=(0, 0, 0), fontname="Helvetica-Bold")
    page.insert_text((72, 140), "A second normal line here.", fontsize=11,
                      color=(0.2, 0.2, 0.7))
    page.draw_rect(fitz.Rect(72, 200, 300, 260), color=(0.8, 0, 0), width=2)
    doc.save(path)
    doc.close()


with tempfile.TemporaryDirectory() as tmp:
    pdf_path = os.path.join(tmp, "sample.pdf")
    docx_path = os.path.join(tmp, "sample.docx")
    pptx_path = os.path.join(tmp, "sample.pptx")
    make_sample_pdf(pdf_path)

    # --- Word ---
    convert_to_docx(pdf_path, docx_path)
    assert os.path.exists(docx_path) and os.path.getsize(docx_path) > 0
    doc = Document(docx_path)
    full_text = "\n".join(p.text for p in doc.paragraphs)
    assert "Hello world" in full_text
    assert "second normal line" in full_text

    # --- PowerPoint ---
    convert_to_pptx(pdf_path, pptx_path)
    assert os.path.exists(pptx_path) and os.path.getsize(pptx_path) > 0
    prs = Presentation(pptx_path)
    assert len(prs.slides) == 1

    shapes = list(prs.slides[0].shapes)
    pictures = [s for s in shapes if s.shape_type == 13]  # PICTURE
    textboxes = [s for s in shapes if s.has_text_frame and s.text_frame.text.strip()]
    assert len(pictures) == 1, "phải có đúng 1 ảnh nền mỗi slide"
    assert len(textboxes) == 2, f"phải có 2 text box (2 dòng chữ gốc), có {len(textboxes)}"

    texts = {tb.text_frame.text for tb in textboxes}
    assert "Hello world, sample paragraph." in texts
    assert "A second normal line here." in texts

    # Giữ đúng định dạng: dòng đầu đậm, dòng hai không đậm và có màu xanh.
    bold_box = next(tb for tb in textboxes if "Hello world" in tb.text_frame.text)
    run = bold_box.text_frame.paragraphs[0].runs[0]
    assert run.font.bold is True

    colored_box = next(tb for tb in textboxes if "second normal line" in tb.text_frame.text)
    run2 = colored_box.text_frame.paragraphs[0].runs[0]
    assert run2.font.bold is False
    assert str(run2.font.color.rgb) == "3333B3"

print("Tất cả self-check pdf_convert đều pass.")


# --- Ảnh nhỏ rời rạc (logo/icon) phải được vẽ RIÊNG thành 1 picture shape
# đúng vị trí/kích thước, không bị khoá cứng vào ảnh nền toàn trang. ---
def make_pdf_with_image(path):
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 100), "Text next to a small logo.", fontsize=14, color=(0, 0, 0))
    img_doc = fitz.open()
    img_page = img_doc.new_page(width=200, height=120)
    img_page.draw_circle((100, 60), 50, color=(1, 0, 0), fill=(1, 0, 0))
    pix = img_page.get_pixmap()
    logo_rect = fitz.Rect(72, 200, 172, 260)  # 100x60pt, khớp tỉ lệ ảnh gốc
    page.insert_image(logo_rect, pixmap=pix)
    doc.save(path)
    doc.close()
    return logo_rect


with tempfile.TemporaryDirectory() as tmp2:
    pdf_path2 = os.path.join(tmp2, "with_image.pdf")
    pptx_path2 = os.path.join(tmp2, "with_image.pptx")
    logo_rect = make_pdf_with_image(pdf_path2)

    convert_to_pptx(pdf_path2, pptx_path2)
    prs2 = Presentation(pptx_path2)
    shapes2 = list(prs2.slides[0].shapes)
    pictures2 = [s for s in shapes2 if s.shape_type == 13]
    assert len(pictures2) == 2, f"phải có ảnh nền + 1 ảnh logo crop riêng, có {len(pictures2)}"

    logo_pic = max(pictures2, key=lambda s: s.left)  # ảnh nền luôn ở (0,0)
    assert logo_pic.left != 0 or logo_pic.top != 0, "ảnh logo phải đặt đúng vị trí gốc, không phải (0,0)"
    got_w = logo_pic.width / EMU_PER_POINT
    got_h = logo_pic.height / EMU_PER_POINT
    assert abs(got_w - logo_rect.width) < 1, (got_w, logo_rect.width)
    assert abs(got_h - logo_rect.height) < 1, (got_h, logo_rect.height)

print("Tất cả self-check ảnh-crop-riêng đều pass.")


# --- Ảnh LỚN (chiếm phần lớn trang, vd ảnh minh họa/biểu đồ full-slide)
# trên trang có chữ thật (KHÔNG phải scan) vẫn phải được crop riêng, không
# bị coi nhầm là "ảnh nền toàn trang" (bug đã report: page_image_rects()
# của translator_engine.py loại ảnh >60% diện tích trang — đúng cho mục
# đích Image Collision Guard lúc dịch, nhưng sai cho mục đích export). ---
def make_pdf_with_large_image(path):
    doc = fitz.open()
    page = doc.new_page(width=600, height=800)
    page.insert_text((20, 20), "Caption above a large illustration.", fontsize=12)
    # PyMuPDF insert_image() giữ nguyên tỉ lệ khung hình (letterbox trong
    # rect chỉ định nếu tỉ lệ không khớp) — dùng đúng tỉ lệ 500:650 của
    # large_rect bên dưới để get_image_rects() trả về ĐÚNG bằng large_rect,
    # không bị co lại theo letterbox.
    img_doc = fitz.open()
    img_page = img_doc.new_page(width=400, height=520)
    img_page.draw_circle((200, 260), 180, color=(0, 0.4, 0.8), fill=(0, 0.4, 0.8))
    pix = img_page.get_pixmap()
    # ~70% diện tích trang -> trước fix bị loại khỏi page_image_rects() (>60%)
    large_rect = fitz.Rect(50, 50, 550, 700)
    page.insert_image(large_rect, pixmap=pix)
    doc.save(path)
    doc.close()
    return large_rect


with tempfile.TemporaryDirectory() as tmp4:
    pdf_path4 = os.path.join(tmp4, "large_image.pdf")
    pptx_path4 = os.path.join(tmp4, "large_image.pptx")
    large_rect = make_pdf_with_large_image(pdf_path4)

    convert_to_pptx(pdf_path4, pptx_path4)
    prs4 = Presentation(pptx_path4)
    pictures4 = [s for s in prs4.slides[0].shapes if s.shape_type == 13]
    assert len(pictures4) == 2, f"ảnh lớn vẫn phải được crop thành shape riêng (nền + ảnh), có {len(pictures4)}"

    big_pic = max(pictures4, key=lambda s: s.width * s.height if (s.left, s.top) != (0, 0) else 0)
    assert (big_pic.left, big_pic.top) != (0, 0), "ảnh lớn phải đặt đúng vị trí gốc, không bị coi là ảnh nền ở (0,0)"
    got_w = big_pic.width / EMU_PER_POINT
    got_h = big_pic.height / EMU_PER_POINT
    assert abs(got_w - large_rect.width) < 1, (got_w, large_rect.width)
    assert abs(got_h - large_rect.height) < 1, (got_h, large_rect.height)

print("Tất cả self-check ảnh-lớn-vẫn-crop đều pass.")


# --- Chữ OCR (trang scan) phải lấy đúng màu mực thật dò được, không phải
# màu đen hardcode của chữ ẩn lúc chèn (bug: _add_text_box từng đọc thẳng
# span["color"], luôn ra đen với mọi trang scan). ---
def make_scanned_pdf_with_red_text(path):
    doc = fitz.open()
    page = doc.new_page(width=400, height=200)
    img_doc = fitz.open()
    img_page = img_doc.new_page(width=400, height=200)
    FONT = "/System/Library/Fonts/Supplemental/Arial.ttf"
    img_page.insert_text((30, 60), "Red scanned heading", fontsize=24,
                          fontfile=FONT, fontname="F1", color=(0.8, 0, 0))
    pix = img_page.get_pixmap(matrix=fitz.Matrix(3, 3))
    page.insert_image(page.rect, pixmap=pix)  # ảnh chiếm toàn trang -> coi là scan
    doc.save(path)
    doc.close()


with tempfile.TemporaryDirectory() as tmp3:
    pdf_path3 = os.path.join(tmp3, "scanned_red.pdf")
    pptx_path3 = os.path.join(tmp3, "scanned_red.pptx")
    make_scanned_pdf_with_red_text(pdf_path3)

    convert_to_pptx(pdf_path3, pptx_path3)
    prs3 = Presentation(pptx_path3)
    shapes3 = list(prs3.slides[0].shapes)
    textboxes3 = [s for s in shapes3 if s.has_text_frame and s.text_frame.text.strip()]
    assert len(textboxes3) == 1, f"phải nhận diện đúng 1 dòng chữ OCR, có {len(textboxes3)}"

    run = textboxes3[0].text_frame.paragraphs[0].runs[0]
    rgb = run.font.color.rgb
    assert rgb != RGBColor(0, 0, 0), "màu chữ OCR không được là đen mặc định (bug đã fix)"
    assert rgb[0] > 150 and rgb[1] < 100 and rgb[2] < 100, f"phải nhận đúng màu mực đỏ thật, ra {rgb}"

print("Tất cả self-check màu-chữ-OCR đều pass.")


# --- _detect_table/_build_table_cells: nhãn hàng wrap 2 dòng ("Foo +" /
# "Bar") phải được gộp thành 1 HÀNG DUY NHẤT (không phải 2 hàng giả), và 1
# ô dữ liệu 3-DÒNG (cao hơn hẳn nhãn hàng NGẮN của chính nó) phải giữ
# TRỌN VẸN trong đúng hàng của nó, không bị dòng cuối "rơi" sang hàng kế
# tiếp dù nằm gần neo hàng sau hơn theo khoảng cách tuyệt đối — 2 bug thật
# đã gặp khi test file PDF thật (bảng OCR từ ảnh scan). Dựng thẳng `lines`
# giả lập (bỏ qua OCR/Vision) để test tất định, không cần ocr_cli."""
def _line(x0, y0, x1, y1, text):
    return ((x0, y0, x1, y1), text, 13.0, (0, 0, 0))


table_lines = [
    _line(10, 10, 60, 25, "Protocol"), _line(100, 10, 160, 25, "Notes"),
    _line(10, 45, 40, 60, "AAA"), _line(100, 45, 220, 60, "Short note one"),
    # Hàng nhãn wrap 2 dòng, gap NHỎ (1pt) — phải gộp thành 1 hàng.
    _line(10, 80, 45, 95, "Foo +"), _line(10, 96, 40, 111, "Bar"),
    # Ô dữ liệu CÙNG hàng, 3 dòng: 2 dòng đầu khớp gap wrap của nhãn, dòng
    # 3 "tiếp tục" xuống thêm (gap 1pt) tới gần sát neo hàng SAU (131).
    _line(100, 80, 260, 95, "Minimizes size and delay"),
    _line(100, 96, 260, 111, "compared with other stuff"),
    _line(100, 112, 260, 127, "final piece of it"),
    _line(10, 131, 40, 146, "BBB"), _line(100, 131, 240, 146, "Another short note"),
]

table_info = _detect_table(table_lines)
assert table_info is not None, "phải nhận diện được bảng giả lập này"
row_y = table_info["row_y"]
assert len(row_y) == 4, f"'Foo +'/'Bar' phải gộp thành 1 hàng — kỳ vọng 4 hàng (header+3), ra {len(row_y)}"

cells = _build_table_cells(table_lines, row_y, table_info["col_x"])
row2_data = [t for _b, t, _s, _c in cells.get((2, 1), [])]
assert row2_data == ["Minimizes size and delay", "compared with other stuff", "final piece of it"], (
    f"cả 3 dòng của ô dữ liệu phải nằm TRỌN trong hàng 2 (Foo+Bar), không rơi dòng nào sang hàng 3, ra {row2_data}"
)
row3_data = [t for _b, t, _s, _c in cells.get((3, 1), [])]
assert row3_data == ["Another short note"], f"hàng 3 (BBB) không được lẫn dòng của hàng 2, ra {row3_data}"

print("Tất cả self-check gộp-nhãn-wrap và gán-hàng-ô-cao đều pass.")


# --- Bảng quá DÀY ĐẶC (ô cần nhiều dòng hơn TABLE_MAX_CELL_LINES, kiểu
# đoạn văn/gạch đầu dòng dài) phải bị TỪ CHỐI dựng thành bảng — bug thật
# đã gặp: 1 dòng LẠC (tiêu đề nhóm, không phải nhãn hàng) tình cờ rơi
# đúng cột nhãn của 1 bảng dày đặc vẫn lọt qua fill_ratio (vẫn lấp đầy),
# kéo lệch toàn bộ hàng phía sau — không tìm được tín hiệu nào tách ca đó
# khỏi bảng thật (bố cục tự do), nên an toàn hơn là từ chối hẳn khi ô quá
# nhiều dòng, để nội dung vẫn vẽ đúng vị trí thật qua text box thường
# thay vì có nguy cơ lệch hàng trong 1 bảng "trông như đúng". ---
dense_table_lines = [
    _line(10, 10, 60, 25, "Protocol"), _line(100, 10, 160, 25, "Notes"),
    _line(10, 45, 40, 60, "AAA"),
    _line(100, 45, 260, 60, "Line one of a very long cell"),
    _line(100, 61, 260, 76, "Line two of a very long cell"),
    _line(100, 77, 260, 92, "Line three of a very long cell"),
    _line(100, 93, 260, 108, "Line four of a very long cell"),
    _line(100, 109, 260, 124, "Line five of a very long cell"),
    _line(10, 150, 40, 165, "BBB"), _line(100, 150, 240, 165, "Short note"),
    _line(10, 181, 40, 196, "CCC"), _line(100, 181, 240, 196, "Another short note"),
]
assert _detect_table(dense_table_lines) is None, "bảng quá dày đặc (ô >4 dòng) phải bị từ chối, không dựng thành bảng"

print("Tất cả self-check từ-chối-bảng-quá-dày-đặc đều pass.")


# --- _page_background_png(table_bbox=...) phải xóa ĐÚNG PHẦN table_bbox
# của 1 nền vector thật (vd hình chữ nhật tô màu xám do PDF vẽ, KHÔNG phải
# ảnh/chữ — _page_background_png cũ không đụng tới loại này), KHÔNG được
# đụng tới phần nền vector còn lại NGOÀI table_bbox — bug thật đã gặp:
# 1 hình chữ nhật nền của bảng gốc thường RỘNG HƠN table_bbox (bảng
# PowerPoint dựng lại chỉ chiếm 1 phần khung layout gốc), redact theo
# graphics=REMOVE_IF_TOUCHED xóa NGUYÊN CẢ hình chữ nhật đó dù chỉ chạm
# 1 phần, làm mất màu nền đúng của phần không thuộc bảng. ---
with tempfile.TemporaryDirectory() as tmp6:
    pdf_path6 = os.path.join(tmp6, "vector_bg.pdf")
    doc6 = fitz.open()
    page6 = doc6.new_page(width=300, height=300)
    page6.draw_rect(fitz.Rect(0, 0, 300, 300), color=None, fill=(0.8, 0.8, 0.8))
    doc6.save(pdf_path6)
    doc6.close()

    bg_doc6 = fitz.open(pdf_path6)
    table_bbox6 = fitz.Rect(100, 100, 200, 200)
    png_bytes6 = _page_background_png(bg_doc6, 0, {}, [], table_bbox=table_bbox6)
    bg_doc6.close()

    pix6 = fitz.Pixmap(png_bytes6)  # ảnh nền render ở zoom 2x (xem _page_background_png)
    inside = pix6.pixel(300, 300)  # điểm (150,150)pt -> giữa table_bbox
    outside = pix6.pixel(100, 100)  # điểm (50,50)pt -> vẫn trong nền xám gốc, NGOÀI table_bbox
    assert inside[:3] == (255, 255, 255), f"bên trong table_bbox phải được xóa trắng, ra {inside}"
    assert outside[:3] != (255, 255, 255), (
        f"bên NGOÀI table_bbox phải GIỮ NGUYÊN màu nền vector gốc (không bị xóa lây), ra {outside}"
    )

print("Tất cả self-check xóa-nền-vector-đúng-vùng-bảng đều pass.")


# --- Kích thước THẬT của bảng PowerPoint (_add_table_shape — tổng độ rộng
# từng cột/tổng chiều cao từng hàng cộng dồn, cái quyết định shape.width/
# height thật sự) phải khớp CHÍNH XÁC với _table_bbox_rect() — bug thật đã
# gặp: col_edges/row_edges lỡ bắt đầu từ col_x[0]/row_y[0] thay vì đúng
# left/top (lệch 4pt, đúng phần đệm left=col_x[0]-4), khiến bảng THẬT nhỏ
# hơn hẳn vùng đã tô trắng nền — để lại viền hở lộ nền/vạch kẻ gốc, nhìn
# như "bảng trong bảng" dù nền đã xóa đúng chỗ. ---
table_info2 = _detect_table(table_lines)
cells2 = _build_table_cells(table_lines, table_info2["row_y"], table_info2["col_x"])
expected_bbox = _table_bbox_rect(table_info2, cells2)

prs2 = Presentation()
slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])
_add_table_shape(slide2, table_info2, cells2)
table_shape = next(s for s in slide2.shapes if s.has_table).table
actual_w = sum(c.width for c in table_shape.columns) / EMU_PER_POINT
actual_h = sum(r.height for r in table_shape.rows) / EMU_PER_POINT

assert abs(actual_w - expected_bbox[2]) < 0.01, (
    f"tổng độ rộng cột thật ({actual_w}) phải khớp _table_bbox_rect() ({expected_bbox[2]}), "
    f"lệch nghĩa là vùng đã xóa nền (theo bbox) và bảng thật vẽ ra không cùng kích thước"
)
assert abs(actual_h - expected_bbox[3]) < 0.01, (
    f"tổng chiều cao hàng thật ({actual_h}) phải khớp _table_bbox_rect() ({expected_bbox[3]})"
)

print("Tất cả self-check kích-thước-bảng-khớp-vùng-xóa-nền đều pass.")


# --- _table_crosses_section_divider: phải phát hiện đúng 1 thanh chia
# nhóm THẬT (gần hết chiều rộng trang, đủ cao — vd nền xám đậm của 1 tiêu
# đề nhóm) cắt ngang qua vùng bảng — bug thật đã gặp: 1 bảng nhỏ vô tình
# trải dài qua 1 thanh như vậy, gộp nhầm nội dung của 2 nhóm khác nhau
# thành 1 bảng. Đồng thời KHÔNG được báo động giả với vạch kẻ mỏng bình
# thường (cao <5pt, vd đường phân cách dòng) nằm trong cùng vùng đó. ---
doc7 = fitz.open()
page7 = doc7.new_page(width=400, height=400)
page7.draw_rect(fitz.Rect(0, 100, 400, 100.5), color=None, fill=(1, 1, 1))  # vạch mỏng, KHÔNG phải chia nhóm
table_bbox7 = fitz.Rect(50, 50, 350, 250)
assert not _table_crosses_section_divider(page7, table_bbox7), "vạch kẻ mỏng bình thường không được coi là thanh chia nhóm"

page7.draw_rect(fitz.Rect(0, 150, 400, 165), color=None, fill=(0.5, 0.5, 0.5))  # thanh chia nhóm THẬT
assert _table_crosses_section_divider(page7, table_bbox7), "phải phát hiện đúng thanh chia nhóm cắt ngang qua vùng bảng"
doc7.close()

print("Tất cả self-check phát-hiện-thanh-chia-nhóm đều pass.")
