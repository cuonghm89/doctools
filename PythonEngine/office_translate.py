"""Dịch trực tiếp file Word (.docx) / PowerPoint (.pptx), giữ nguyên định
dạng gốc — khác hẳn cách translator_engine.py xử lý PDF (phải xóa/vẽ lại
từng khối chữ bằng redaction, vì PDF không có khái niệm "run chữ sửa
được"): .docx/.pptx đã là định dạng có cấu trúc chữ sẵn (run/paragraph),
chỉ cần THAY TRỰC TIẾP nội dung chữ trong các run đó — mọi định dạng khác
(font, cỡ, màu, in đậm, vị trí, bảng biểu, ảnh...) tự động giữ nguyên vì
không đụng tới XML nào khác.

1 đoạn văn có thể bị chia thành NHIỀU run (do đổi định dạng giữa chừng câu,
hoặc do chính Word/PowerPoint tự tách run mà không có lý do rõ ràng) — dịch
từng run riêng lẻ sẽ cho ra bản dịch vô nghĩa (thiếu ngữ cảnh, đứt câu giữa
chừng). Nên: gộp toàn bộ text của 1 đoạn văn thành 1 đơn vị dịch, dịch cả
đoạn cùng lúc, rồi ghi hết bản dịch vào RUN ĐẦU TIÊN (giữ định dạng của
run đó cho toàn bộ đoạn), xóa nội dung mọi run còn lại trong cùng đoạn.
ponytail: mất định dạng khác nhau GIỮA các run trong cùng 1 đoạn (vd nửa
câu in đậm, nửa không) — chấp nhận được vì hiếm và không ảnh hưởng bố cục
tổng thể; nâng cấp lên giữ đúng từng run cần tự chia lại bản dịch theo vị
trí ký tự tương ứng, phức tạp hơn nhiều so với lợi ích thực tế.

Khác với PDF, bảng trong .docx/.pptx là bảng THẬT (cấu trúc row/cell rõ
ràng, không phải suy đoán từ vị trí chữ như _detect_table() trong
pdf_convert.py) — không cần heuristic is_table_lines() nào để ĐỊNH TUYẾN
dịch (mọi đoạn văn, dù trong ô bảng hay không, đều gửi cho router với
is_table=False như nhau). Cờ `in_table` dùng ở file này mang Ý NGHĨA KHÁC:
chỉ để CHỌN ngưỡng co chữ tối thiểu (ô bảng chật hơn text box tự do, chấp
nhận co nhiều hơn để đỡ tràn cột/hàng) — xem FONT_SHRINK_MIN_SCALE_TABLE.
"""
import os

import fitz  # PyMuPDF — đo độ rộng chữ thật, xem _text_width_ratio()
from docx import Document
from docx.shared import Pt as DocxPt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu, Pt as PptxPt

from translator_engine import emit, VN_FONT_DIR

# Tiếng Việt dài hơn tiếng Anh (~20-30%, có khi hơn với câu ngắn) — khác
# PDF (khung pixel cố định, fit_and_draw() tự co/giãn/cắt), .docx/.pptx
# không có khung pixel để đo trực tiếp — co cỡ chữ theo tỉ lệ ĐỘ RỘNG THẬT
# phình ra (_text_width_ratio()), chặn dưới ở FONT_SHRINK_MIN_SCALE* để
# không co tới mức không đọc được. Ô bảng chật hơn hẳn text box tự do
# (không có chỗ giãn dòng thoải mái, cột hẹp cố định) nên chấp nhận co
# nhiều hơn — ngưỡng riêng thấp hơn.
FONT_SHRINK_MIN_SCALE = 0.55
FONT_SHRINK_MIN_SCALE_TABLE = 0.40

# Co chữ tới ngưỡng vẫn có thể không đủ (bản dịch phình quá nhiều) — với ô
# bảng, lớp phòng thủ THỨ 2 là GIÃN CHIỀU CAO HÀNG cho đủ chỗ (chỉ giãn,
# không co nhỏ lại), đo số dòng cần bằng word-wrap thật trên font thật thay
# vì đoán — đây là cách duy nhất ĐẢM BẢO hết tràn, khác co chữ (có chặn
# dưới, không đảm bảo tuyệt đối).
EMU_PER_POINT = 12700
LINE_HEIGHT_FACTOR = 1.2  # hệ số giãn dòng thông thường trên cỡ chữ (typography phổ biến)
DEFAULT_CELL_MARGIN_LR_EMU = 91440   # ~0.1in — margin trái/phải mặc định của ô bảng PowerPoint khi chưa set riêng
DEFAULT_CELL_MARGIN_TB_EMU = 45720   # ~0.05in — margin trên/dưới mặc định

_METRIC_FONT = fitz.Font(fontfile=os.path.join(VN_FONT_DIR, "Arial.ttf"))
_METRIC_FONT_SIZE = 100  # cỡ đo cố định — tỉ lệ độ rộng tuyến tính theo fontsize nên cỡ đo không ảnh hưởng kết quả


def _wrapped_line_count(text, font_size_pt, max_width_pt):
    """Số dòng cần để word-wrap `text` vừa trong `max_width_pt`, đo bằng
    chính font Arial thật ở `font_size_pt` (fitz.Font.text_length) — thuật
    toán word-wrap THAM LAM (greedy): xếp từ vào dòng hiện tại tới khi
    không vừa thì xuống dòng mới, giống cách trình soạn thảo thật ngắt
    dòng. Chính xác hơn hẳn ước lượng, vì đo đúng độ rộng từng từ thật."""
    words = text.split()
    if not words or max_width_pt <= 0:
        return 1
    space_w = _METRIC_FONT.text_length(" ", fontsize=font_size_pt)
    lines = 1
    current_w = 0.0
    for word in words:
        word_w = _METRIC_FONT.text_length(word, fontsize=font_size_pt)
        needed = word_w if current_w == 0 else current_w + space_w + word_w
        if needed > max_width_pt and current_w > 0:
            lines += 1
            current_w = word_w
        else:
            current_w = needed
    return lines


def _grow_table_rows_to_fit(table):
    """Sau khi dịch+co chữ xong (_translate_paragraphs), GIÃN chiều cao
    từng hàng nếu nội dung (đã co chữ, đo bằng word-wrap thật) vẫn cần
    nhiều dòng hơn hàng gốc chứa nổi. Chỉ GIÃN, không bao giờ co nhỏ lại —
    giữ nguyên tỉ lệ trình bày gốc khi bản dịch không cần thêm chỗ."""
    columns = table.columns
    for row in table.rows:
        needed_height_emu = 0
        for col_idx, cell in enumerate(row.cells):
            text = cell.text_frame.text
            if not text.strip():
                continue
            runs = [run for p in cell.text_frame.paragraphs for run in p.runs if run.text.strip()]
            font_pt = max((run.font.size.pt for run in runs if run.font.size), default=12.0)
            margin_l = cell.margin_left if cell.margin_left is not None else DEFAULT_CELL_MARGIN_LR_EMU
            margin_r = cell.margin_right if cell.margin_right is not None else DEFAULT_CELL_MARGIN_LR_EMU
            usable_width_pt = (columns[col_idx].width - margin_l - margin_r) / EMU_PER_POINT
            lines = _wrapped_line_count(text, font_pt, usable_width_pt)
            margin_t = cell.margin_top if cell.margin_top is not None else DEFAULT_CELL_MARGIN_TB_EMU
            margin_b = cell.margin_bottom if cell.margin_bottom is not None else DEFAULT_CELL_MARGIN_TB_EMU
            cell_height_pt = lines * font_pt * LINE_HEIGHT_FACTOR + (margin_t + margin_b) / EMU_PER_POINT
            needed_height_emu = max(needed_height_emu, int(cell_height_pt * EMU_PER_POINT))
        if needed_height_emu > row.height:
            row.height = Emu(needed_height_emu)


def _iter_pptx_tables(shapes):
    """Đệ quy vào group shape, trả về mọi bảng thật (table object) trên 1
    tập shape. Dùng để giãn hàng sau khi dịch — xem _grow_table_rows_to_fit()."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pptx_tables(shape.shapes)
        elif shape.has_table:
            yield shape.table


def _text_width_ratio(original, translated):
    """Tỉ lệ ĐỘ RỘNG THẬT (đo bằng chính font Arial dùng để vẽ tiếng Việt
    trong translator_engine.py, qua fitz.Font.text_length) giữa bản dịch
    và bản gốc — chính xác hơn hẳn đếm số ký tự (vd "iiii" hẹp hơn nhiều
    "wwww" dù cùng 4 ký tự, chữ có dấu tiếng Việt cũng rộng hơn chữ
    không dấu). Chỉ đo độ rộng 1 DÒNG, không mô phỏng lại việc xuống dòng
    nhiều dòng — vẫn là xấp xỉ, không chính xác tuyệt đối như
    fit_and_draw() (PDF) tự thử vẽ thật rồi mới biết vừa hay không."""
    orig_w = _METRIC_FONT.text_length(original, fontsize=_METRIC_FONT_SIZE)
    if orig_w <= 0:
        return 1.0
    return _METRIC_FONT.text_length(translated, fontsize=_METRIC_FONT_SIZE) / orig_w


def _apply_uppercase(original, translated):
    # Cùng lý do đã áp dụng trong process_pdf(): DeepL/Gemini không giữ
    # nguyên chữ HOA gốc, phải tự ép lại nếu bản gốc toàn chữ hoa.
    return translated.upper() if original.strip() and original.isupper() else translated


def _paragraph_text(paragraph):
    return "".join(run.text for run in paragraph.runs)


def _adjust_font_for_length(run, original, translated, pt_class, min_scale):
    """Co cỡ chữ của run tỉ lệ nghịch với mức phình dài (đo bằng
    _text_width_ratio()), để giảm tràn/chồng lấn khi bản dịch dài hơn hẳn
    bản gốc. Chỉ áp dụng khi run ĐÃ set cỡ chữ riêng (run.font.size không
    None) — nếu cỡ chữ kế thừa từ placeholder/theme (None), không có cơ sở
    để biết đang co từ bao nhiêu, an toàn hơn là bỏ qua thay vì đoán bừa."""
    if run.font.size is None or not original.strip():
        return
    ratio = _text_width_ratio(original, translated)
    if ratio <= 1.05:
        return
    scale = max(min_scale, 1 / ratio)
    run.font.size = pt_class(run.font.size.pt * scale)


def _write_paragraph(paragraph, original, translated, pt_class, min_scale):
    runs = paragraph.runs
    if not runs:
        return
    runs[0].text = translated
    _adjust_font_for_length(runs[0], original, translated, pt_class, min_scale)
    for run in runs[1:]:
        run.text = ""


def _translate_paragraphs(paragraph_specs, router, page, total, pt_class):
    """Dùng chung cho cả docx và pptx. `paragraph_specs`: list (paragraph,
    in_table) — xem docstring đầu file về ý nghĩa `in_table` ở đây (chọn
    ngưỡng co chữ, KHÔNG phải để định tuyến DeepL/Gemini). `page`/`total`
    chỉ để báo tiến trình (số đoạn văn đã xử lý cho docx, số slide cho
    pptx — xem 2 hàm bên dưới), không mang ý nghĩa "trang" thật với docx.

    enforce_length_guard=False: khác PDF (khung cố định, guard "phình quá
    dài thì giữ bản gốc" là lựa chọn ĐÚNG vì không có chỗ tự co giãn),
    docx/pptx tự co cỡ chữ (_adjust_font_for_length) thay vì bỏ dịch — bật
    guard này ở đây sẽ tạo ra tiêu đề/đoạn văn nửa Anh nửa Việt khó hiểu."""
    texts = [_paragraph_text(p) for p, _in_table in paragraph_specs]
    unit_indices = [i for i, t in enumerate(texts) if t.strip()]
    units = [(texts[i], False) for i in unit_indices]
    results = router.translate_batch(units, enforce_length_guard=False) if units else []

    for para_idx, result in zip(unit_indices, results):
        translated, engine, deepl_error, item_error = result
        if item_error is not None:
            emit(type="progress", page=page, total=total, engine="skipped", detail=item_error)
            continue
        original = texts[para_idx]
        paragraph, in_table = paragraph_specs[para_idx]
        min_scale = FONT_SHRINK_MIN_SCALE_TABLE if in_table else FONT_SHRINK_MIN_SCALE
        _write_paragraph(paragraph, original, _apply_uppercase(original, translated), pt_class, min_scale)
        emit(type="progress", page=page, total=total, engine=engine, detail=deepl_error)


def _iter_docx_table_paragraphs(table):
    for row in table.rows:
        for cell in row.cells:
            for p in cell.paragraphs:
                yield p, True
            for nested in cell.tables:
                yield from _iter_docx_table_paragraphs(nested)


def _iter_docx_paragraphs(doc):
    """Mọi đoạn văn cần dịch: nội dung chính (body), bảng (đệ quy — 1 ô có
    thể chứa bảng con), rồi header/footer của mỗi section. Trả về
    (paragraph, in_table)."""
    for p in doc.paragraphs:
        yield p, False
    for table in doc.tables:
        yield from _iter_docx_table_paragraphs(table)
    for section in doc.sections:
        for part in (section.header, section.footer):
            for p in part.paragraphs:
                yield p, False
            for table in part.tables:
                yield from _iter_docx_table_paragraphs(table)


def translate_docx(input_path, output_path, router, max_pages=0):
    """max_pages bị bỏ qua: .docx không có khái niệm "trang" cố định trong
    XML (phụ thuộc render/pagination lúc mở bằng Word thật), không có cách
    nào đáng tin cậy để chỉ dịch N trang đầu."""
    doc = Document(input_path)
    paragraph_specs = list(_iter_docx_paragraphs(doc))
    total = len(paragraph_specs)
    _translate_paragraphs(paragraph_specs, router, page=total, total=total, pt_class=DocxPt)
    doc.save(output_path)


def _iter_pptx_text_frames(shapes):
    """Mọi text_frame cần dịch trong 1 tập shape (đệ quy vào group shape và
    bảng — 1 GraphicFrame có bảng thì text nằm trong từng cell, không phải
    ở chính shape đó). Trả về (text_frame, is_table_cell)."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pptx_text_frames(shape.shapes)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame, True
        elif shape.has_text_frame:
            yield shape.text_frame, False


def translate_pptx(input_path, output_path, router, max_pages=0):
    prs = Presentation(input_path)
    slides = list(prs.slides)
    total = min(max_pages, len(slides)) if max_pages else len(slides)

    for slide_index, slide in enumerate(slides):
        if slide_index >= total:
            break
        text_frames = list(_iter_pptx_text_frames(slide.shapes))
        for tf, is_table_cell in text_frames:
            # word_wrap: bắt buộc xuống dòng thay vì tràn NGANG sang cột/
            # shape bên cạnh khi bản dịch dài hơn — áp cho cả text box lẫn
            # ô bảng, kể cả khi file gốc lỡ tắt wrap.
            tf.word_wrap = True
            if not is_table_cell:
                # Bật tính năng tự co chữ CÓ SẴN của PowerPoint (normAutofit)
                # cho text box thường — PowerPoint thật tự tính lại
                # fontScale khi mở file nếu chữ vẫn tràn khung, chính xác
                # hơn hẳn _adjust_font_for_length() vì đo bằng font engine
                # thật của nó, không phải ước lượng. KHÔNG bật cho ô bảng:
                # PowerPoint tự quản lý layout bảng riêng, normAutofit trên
                # 1 cell không có hành vi đáng tin cậy tương đương text box.
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        paragraph_specs = [
            (para, is_table_cell)
            for tf, is_table_cell in text_frames
            for para in tf.paragraphs
        ]
        _translate_paragraphs(paragraph_specs, router, page=slide_index + 1, total=total, pt_class=PptxPt)

        for table in _iter_pptx_tables(slide.shapes):
            _grow_table_rows_to_fit(table)

    prs.save(output_path)
