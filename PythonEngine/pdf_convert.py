"""Chuyển đổi PDF -> Word (.docx) / PowerPoint (.pptx). Độc lập với pipeline
dịch thuật: dùng được cho bất kỳ PDF nào (bản gốc hoặc bản đã dịch).

Word: giao cho pdf2docx, một thư viện chuyên dụng đã được kiểm chứng rộng
rãi — không có lý do gì để tự viết lại việc tái tạo PDF thành văn bản dạng
dòng chảy (nhận diện đoạn văn/bảng biểu) khi đã có thư viện làm tốt việc đó.

PowerPoint: chưa có thư viện "PDF->PPTX chỉnh sửa được" nào đủ trưởng thành
(hầu hết chỉ dán ảnh chụp nguyên trang vào mỗi slide). Tự xây ở đây, tái sử
dụng kỹ thuật đã được kiểm chứng trong translator_engine.py:
  - Lớp nền: xóa (redact) CHỈ phần chữ (ảnh/vector giữ nguyên, dùng đúng
    PDF_REDACT_IMAGE_NONE / graphics=0 như bên translator) trên một bản sao
    dùng-rồi-bỏ của tài liệu, sau đó render mỗi trang thành ảnh nền cho slide.
  - Lớp chữ: text box PowerPoint thật, sửa được, đặt đúng vị trí/kích
    thước/màu/đậm-nghiêng của từng khối chữ gốc, trích từ tài liệu GỐC
    (chưa bị đụng tới).
Cách này giữ ảnh/vector chính xác từng pixel (dưới dạng ảnh nền) trong khi
chữ vẫn sửa được thật sự trong PowerPoint — đánh đổi lấy rủi ro nhỏ về việc
vị trí/ngắt dòng của text box có thể lệch chút (engine dàn chữ của
PowerPoint khác của PDF), đúng theo lựa chọn bạn đã chọn.

Trang PDF dạng scan (không có lớp chữ) được OCR qua ocr_pdf.py trước khi
trích khối chữ, nên vẫn tạo được text box sửa được như trang số hóa bình
thường. Ảnh nền của trang scan cũng được xóa đúng pixel vùng chữ gốc (dùng
rect Vision nhận diện được, PDF_REDACT_IMAGE_PIXELS) trước khi render —
nếu không, chữ scan gốc sẽ lộ ra dưới text box mới chồng lên trên.
"""
import io
import logging
import math
import os
import tempfile

import fitz  # PyMuPDF

from ocr_pdf import ensure_text_layer
from paragraphs import merge_paragraph_blocks, split_incoherent_block
from translator_engine import local_bg_color

# pdf2docx (bản mới nhất, 0.5.8, đã ngừng cập nhật) gọi Rect.get_area() —
# hàm này đã bị PyMuPDF gỡ bỏ ở các bản mới (ta cũng từng dính lỗi này với
# code của chính mình). Không thể hạ cấp PyMuPDF vì engine dịch đang cần
# bản mới cho insert_htmlbox, nên vá tương thích ngược 1 dòng thay vì tách
# venv riêng cho mỗi thư viện.
if not hasattr(fitz.Rect, "get_area"):
    fitz.Rect.get_area = lambda self: self.width * self.height

from pdf2docx import Converter

# pdf2docx gọi thẳng logging.basicConfig(level=INFO) lên ROOT logger ngay
# khi import (không dùng logger riêng theo tên module), và in log ra
# stdout. App gọi engine này qua subprocess và đọc stdout từng dòng như
# JSON (xem translator_engine.py) — log lẫn vào sẽ phá giao thức đó, nên
# phải hạ cấp độ root logger XUỐNG SAU KHI import xong.
logging.getLogger().setLevel(logging.CRITICAL)
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

EMU_PER_POINT = 12700  # python-pptx tính bằng EMU; tọa độ PDF tính bằng point


def convert_to_docx(input_pdf, output_docx):
    # pdf2docx chỉ nhận đường dẫn file, không nhận Document đang mở trong
    # bộ nhớ — nên OCR trước, lưu ra 1 file tạm đã có lớp chữ, rồi mới đưa
    # file tạm đó cho pdf2docx xử lý như bình thường.
    doc = fitz.open(input_pdf)
    # visible=True: pdf2docx tự lọc bỏ hoàn toàn chữ vô hình (render_mode=3)
    # khi trích xuất (đã kiểm chứng thực tế) — file tạm này chỉ pdf2docx
    # đọc, người dùng không bao giờ mở nó dưới dạng PDF, nên chữ "hiện" ở
    # đây an toàn, không tạo ra trùng lặp hình ảnh nào cho người dùng thấy.
    ensure_text_layer(doc, visible=True)
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp_path = tmp.name
    try:
        doc.save(tmp_path)
        doc.close()
        cv = Converter(tmp_path)
        try:
            cv.convert(output_docx)
        finally:
            cv.close()
    finally:
        os.unlink(tmp_path)


def _page_background_png(bg_doc, page_index, ocr_rects_by_page, image_rects, table_bbox=None):
    """Render 1 trang sau khi đã xóa hết chữ VÀ các ảnh nhỏ rời rạc (ảnh
    được vẽ lại riêng, thành picture shape sửa được — xem
    `_crop_image_png()`/`convert_to_pptx()`) — chỉ còn nền màu/vector.

    `table_bbox`: nếu trang này SẼ được dựng thành bảng PowerPoint thật
    (`_add_table_shape()`), phải xóa luôn phần NỀN/VẠCH KẺ VECTOR THẬT của
    bảng gốc trong vùng đó (bug thật đã gặp: bảng gốc dùng hình chữ nhật tô
    màu xám + vạch kẻ trắng vẽ bằng vector, KHÔNG phải ảnh/chữ, nên lượt
    xóa chữ ở trên không đụng tới — bảng PowerPoint dựng lại đè lên trên
    một bảng gốc vẫn còn nguyên bên dưới, không khớp pixel-perfect, nhìn
    như "2 bảng chồng nhau", rõ nhất ở góc lệch nhiều nhất).

    Trang scan (toàn bộ trang là 1 ảnh) xử lý khác hẳn trang số hóa bình
    thường: `bg_doc` không được OCR (đây là bản dùng-rồi-bỏ, chỉ để render
    ảnh nền) nên get_text("dict") không thấy gì để xóa — nếu không xử lý
    riêng, chữ gốc trong ảnh sẽ giữ nguyên, còn text box sửa-được mới lại
    chồng thẳng lên trên, ra 2 lớp chữ đè nhau. Dùng thẳng rect GỐC do
    Vision nhận diện (đã có sẵn từ lúc OCR `doc` ở convert_to_pptx) để xóa
    đúng pixel chữ cũ — không đi qua get_text("dict")/gộp khối nào, vì
    PyMuPDF có thể tự gộp nhiều dòng KHÔNG liên quan thành 1 "block" (xem
    translator_engine.py), làm vùng xóa lem sang chỗ không phải chữ."""
    page = bg_doc[page_index]
    is_ocr_page = page_index in ocr_rects_by_page
    if is_ocr_page:
        text_rects = [rect for rect, _color in ocr_rects_by_page[page_index]]
    else:
        blocks = page.get_text("dict")["blocks"]
        text_rects = [fitz.Rect(b["bbox"]) for b in blocks if b.get("type") == 0]
    # Lấy mẫu màu nền quanh từng vùng chữ TRƯỚC khi xóa (trang vẫn còn
    # nguyên vẹn) — tô lại đúng màu nền thật thay vì trắng phẳng mặc định,
    # nếu không sẽ để lại các mảng trắng lạc quẻ trên nền màu/gradient.
    bg_colors = [local_bg_color(page, rect) for rect in text_rects]
    for rect, bg_color in zip(text_rects, bg_colors):
        page.add_redact_annot(rect, fill=bg_color)
    if text_rects:
        images_mode = fitz.PDF_REDACT_IMAGE_PIXELS if is_ocr_page else fitz.PDF_REDACT_IMAGE_NONE
        page.apply_redactions(images=images_mode, graphics=0)

    # Ảnh rời rạc (logo/icon/biểu đồ/ảnh minh họa — xem
    # `_page_image_rects_for_export()`) được xóa ở đây bằng 1 lượt
    # apply_redactions() RIÊNG, dùng images=PIXELS: tách riêng khỏi lượt
    # xóa chữ ở trên để không nới lỏng an toàn của nó (trang số hóa cố ý
    # dùng images=NONE, không bao giờ cắt lẹm vào ảnh dù 1 khối chữ tình cờ
    # đè nhẹ lên rìa ảnh).
    if image_rects:
        for rect in image_rects:
            page.add_redact_annot(rect, fill=local_bg_color(page, rect))
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS, graphics=0)

    # Vùng bảng (nếu có): VẼ ĐÈ 1 hình chữ nhật trắng lên trên — KHÔNG dùng
    # add_redact_annot()/apply_redactions(graphics=...) như 2 lượt trên.
    # Đã thử graphics=PDF_REDACT_LINE_ART_REMOVE_IF_TOUCHED trước đó và
    # gặp 2 bug: (1) nó xóa TOÀN BỘ path vector "chạm" vào rect, không chỉ
    # phần NẰM TRONG rect — 1 hình chữ nhật nền xám của bảng gốc thường
    # rộng hơn hẳn table_bbox (kéo dài sang cả phần mô tả bên trái không
    # thuộc bảng), nên bị xóa mất màu nền đúng của phần đó, thành tác dụng
    # phụ ngoài ý muốn; (2) fill lấy qua local_bg_color(page, table_bbox)
    # đọc lại CHÍNH vùng đang xóa SAU KHI chữ đã bị xóa nhưng nền xám vector
    # thì CHƯA — vô tình đọc lại đúng màu xám đó rồi tô y hệt, coi như
    # không đổi gì. Vẽ đè trực tiếp (không phải redact) tránh cả 2: chỉ phủ
    # ĐÚNG PHẦN diện tích table_bbox bằng trắng cố định (khớp nền trắng/
    # theme mặc định của bảng PowerPoint dựng lại — _add_table_shape),
    # không đụng gì tới phần còn lại của trang.
    if table_bbox is not None:
        page.draw_rect(table_bbox, color=None, fill=(1, 1, 1))

    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # phóng 2x cho ảnh nền nét hơn
    return pix.tobytes("png")


def _page_image_rects_for_export(page, is_ocr_page):
    """Mọi ảnh nhúng rời rạc trên trang, để crop RIÊNG thành picture shape
    sửa được trong convert_to_pptx(). KHÔNG dùng
    translator_engine.page_image_rects(): hàm đó loại ảnh chiếm >60% diện
    tích trang, nhưng ngưỡng đó dành cho mục đích khác hẳn (Image Collision
    Guard lúc DỊCH — chỉ cần bảo vệ logo/icon nhỏ khỏi bị đè chữ dịch lên);
    áp ngưỡng đó vào đây khiến mọi ảnh lớn (ảnh minh họa, biểu đồ chiếm gần
    hết trang — rất phổ biến trong slide/tài liệu thật) bị coi nhầm là
    "ảnh nền toàn trang" và không được crop, đúng thứ người dùng report.

    Trang scan (is_ocr_page — không có chữ thật, `ensure_text_layer()` đã
    OCR) thì KHÔNG có ảnh "rời rạc" nào để crop: ảnh CHÍNH LÀ nội dung cả
    trang, giữ nguyên trong ảnh nền như translator_engine.py vẫn làm."""
    if is_ocr_page:
        return []
    rects = []
    for img in page.get_images(full=True):
        xref = img[0]
        rects.extend(page.get_image_rects(xref))
    return rects


def _crop_image_png(page, rect, zoom=2.0):
    """Render đúng 1 vùng ảnh nhỏ (logo/icon/biểu đồ) từ trang GỐC (chưa bị
    `_page_background_png()` xóa), để chèn thành 1 picture shape riêng, sửa
    được (di chuyển/đổi cỡ/xóa độc lập) thay vì bị khoá cứng vào ảnh nền
    toàn trang như trước.
    ponytail: render lại bằng get_pixmap (phẳng, không giữ alpha) thay vì
    trích byte ảnh gốc qua extract_image() — đơn giản hơn nhiều (không phải
    lo colorspace/mask của ảnh nhúng), khớp pixel-perfect với những gì thấy
    trên trang, nhưng ảnh có nền trong suốt sẽ mất độ trong suốt đó (không
    đáng lo vì phía dưới đã được tô đúng màu nền thật ở _page_background_png)."""
    pix = page.get_pixmap(clip=rect, matrix=fitz.Matrix(zoom, zoom))
    return pix.tobytes("png")


def _add_text_box(slide, group, page_height, ocr_rects=None, scale=1.0):
    """`group` là 1 nhóm đã gộp từ merge_paragraph_blocks (có thể gồm
    nhiều "sub_blocks" — các dòng wrap của cùng 1 đoạn văn/bullet). Dùng
    CHUNG 1 cỡ chữ (group["size"], trung bình cỡ chữ thật của cả đoạn) cho
    mọi dòng trong nhóm — không lấy riêng cỡ chữ ước lượng của từng dòng
    (span["size"]): với chữ OCR, mỗi dòng được ước lượng cỡ chữ ĐỘC LẬP
    với nhau, sai số nhỏ giữa các dòng khiến dòng đầu bullet và dòng xuống
    dòng tiếp theo hiện cỡ chữ khác nhau rõ rệt dù trong ảnh gốc chúng
    hoàn toàn cùng 1 cỡ.

    Nhãn xoay dọc (vd "Applications"/"Communications" cạnh sơ đồ OSI) được
    ocr_pdf.py chèn với rotate=90 — PyMuPDF ghi lại hướng đọc thật của dòng
    đó vào line["dir"] (khác (1,0) mặc định). Text box PowerPoint phải
    xoay CẢ KHỐI (thuộc tính .rotation, độ, chiều kim đồng hồ) để khớp:
    tạo box ở kích thước HOÁN ĐỔI rộng/cao (vì PowerPoint xoay quanh tâm
    của box được tạo ra chưa xoay), đặt tâm trùng với tâm vùng đích, rồi
    mới xoay — nếu không, box sẽ xoay lệch tâm ra khỏi đúng vị trí trực
    quan mong muốn."""
    x0, y0, x1, y1 = group["bbox"]
    width = max(x1 - x0, 1)
    height = max(y1 - y0, 1)

    first_line = group["sub_blocks"][0]["lines"][0]
    dx, dy = first_line.get("dir", (1.0, 0.0))
    is_rotated = abs(dx) < 0.5
    box_w, box_h = (height, width) if is_rotated else (width, height)
    cx, cy = x0 + width / 2, y0 + height / 2
    left, top = cx - box_w / 2, cy - box_h / 2

    box = slide.shapes.add_textbox(
        Emu(int(left * scale * EMU_PER_POINT)),
        Emu(int(top * scale * EMU_PER_POINT)),
        Emu(int(box_w * scale * EMU_PER_POINT)),
        Emu(int(box_h * scale * EMU_PER_POINT)),
    )
    if is_rotated:
        box.rotation = math.degrees(math.atan2(dy, dx))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    font_size = Pt(max(group["size"], 1))
    first_para = True
    for block in group["sub_blocks"]:
        for line in block["lines"]:
            para = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            # Chữ OCR (vô hình) luôn mang màu đen hardcode ở span["color"]
            # lúc chèn (không biết màu thật) — dùng _line_color() để khớp
            # với màu mực THẬT dò được từ ảnh scan (giống hệt cách bảng
            # biểu đã làm đúng ở _add_table_shape/_page_lines); trang không
            # OCR thì _line_color() tự rơi về đọc span["color"] như cũ.
            r, g, b = _line_color(line, ocr_rects)
            for span in line["spans"]:
                run = para.add_run()
                run.text = span["text"]
                run.font.size = font_size
                flags = span.get("flags", 0)
                run.font.bold = bool(flags & 16)
                run.font.italic = bool(flags & 2)
                run.font.color.rgb = RGBColor(r, g, b)
    return box


# Bảng biểu OCR (vd bảng so sánh thông số) trước đây chỉ được vẽ bằng ảnh
# nền (giữ đúng đường kẻ gốc) + các text box rời rạc đặt đè lên — mỗi ô là
# 1 dòng nhận diện riêng của Vision nên vị trí dễ lệch khỏi đúng tâm ô,
# nhìn lộn xộn. Các hàm dưới đây dựng lại thành 1 bảng PowerPoint THẬT
# (table object, sửa được, căn ô chuẩn) thay vì ảnh + chữ đè.
#
# ponytail: trang OCR là 1 ảnh raster, không có dữ liệu đường kẻ vector để
# đọc lại chính xác — số hàng/cột và NỘI DUNG từng ô suy ra khá đáng tin
# cậy từ vị trí chữ (xem bên dưới), nhưng ĐỘ RỘNG hàng/cột cuối cùng chỉ
# là ước lượng theo khoảng cách giữa các điểm neo phát hiện được, không
# sao chép chính xác tỉ lệ ảnh gốc.
MIN_TABLE_ROWS = 3
MIN_TABLE_COLS = 2
TABLE_COL_GAP = 25
TABLE_ROW_LABEL_SIZE_TOLERANCE = 0.5
TABLE_LABEL_MAX_MEDIAN_LEN = 25
MAX_LABEL_TO_DATA_RATIO = 2.0
TABLE_MIN_FILL_RATIO = 0.9
# Từ chối dựng bảng nếu bất kỳ ô nào cần nhiều dòng hơn ngưỡng này — xem
# giải thích đầy đủ tại chỗ dùng trong _detect_table().
TABLE_MAX_CELL_LINES = 4
# Bù lệch canh dòng cơ sở giữa cột nhãn/cột dữ liệu khi gán dòng vào hàng
# (_interval_index, chỉ dùng cho trục hàng) — CỐ ĐỊNH, không suy ra từ
# khoảng cách hàng trung bình (dễ bị 1 hàng cao bất thường kéo lệch, xem
# _build_table_cells()).
TABLE_ROW_BASELINE_TOLERANCE = 5.0


def _cluster_1d(items, gap):
    """items: list (vị_trí, payload). Gom thành các dải liên tục — bắt đầu
    dải mới khi khoảng cách với phần tử ngay trước (đã sắp xếp) vượt quá
    `gap`. Trả về list các dải, mỗi dải là list (vị_trí, payload)."""
    if not items:
        return []
    ordered = sorted(items, key=lambda t: t[0])
    bands = [[ordered[0]]]
    for item in ordered[1:]:
        if item[0] - bands[-1][-1][0] > gap:
            bands.append([])
        bands[-1].append(item)
    return bands


def _label_continues_in_other_columns(lines, y0, y1, label_col_x_max, tolerance=3.0):
    """Có dòng nào ở CÁC CỘT KHÁC (không phải cột nhãn) bắt đầu trong
    khoảng [y0-tolerance, y1+tolerance] và mở đầu bằng CHỮ THƯỜNG không —
    dấu hiệu 1 câu đang viết dở tiếp tục xuống dòng (cùng 1 Ô, chỉ là nhãn
    hàng hẹp bên cạnh tình cờ cũng wrap xuống đúng lúc đó), khác hẳn 1 ý/
    câu MỚI bắt đầu (viết hoa) của 1 HÀNG THẬT KHÁC. Dùng làm tín hiệu phụ
    cho _merge_wrapped_row_labels() — riêng khoảng cách dọc giữa 2 dòng
    nhãn không đủ để phân biệt 2 ca này (đã đo thực tế: nhãn wrap thật
    "CoAP +"/"CBOR" cách nhau 0.3pt, nhưng 2 NHÃN GIỐNG HỆT của 2 hàng
    khác nhau "WebSocket"/"WebSocket" cũng chỉ cách nhau 1.2pt — cùng
    thang độ lớn, không thể tách bằng ngưỡng khoảng cách đơn thuần)."""
    for bbox, text, _size, _color in lines:
        if bbox[0] < label_col_x_max:
            continue
        if y0 - tolerance <= bbox[1] <= y1 + tolerance:
            stripped = text.strip()
            if stripped and stripped[0].islower():
                return True
    return False


def _merge_wrapped_row_labels(bboxes_sorted_by_y, lines, label_col_x_max, gap_factor=1.0):
    """1 nhãn hàng (label) có thể tự xuống dòng (vd "CoAP + CBOR" hiện
    thành 2 dòng OCR riêng "CoAP +"/"CBOR") — mỗi dòng NHƯ VẬY vẫn chỉ là
    1 HÀNG THẬT duy nhất, không phải 2 hàng; nếu không gộp lại, dò neo hàng
    sẽ tách nhầm thành 2 hàng giả, kéo theo mọi nội dung cột bên bị gán
    lệch hàng từ đây trở đi. Gộp 2 dòng liên tiếp làm 1 khi khoảng cách dọc
    giữa chúng NHỎ (dưới gap_factor lần chiều cao dòng trước — loại được
    phần lớn trường hợp rõ ràng là 2 hàng khác nhau, vốn cách nhau xa hơn
    hẳn) VÀ có tín hiệu phụ xác nhận đang viết dở câu ở cột bên
    (_label_continues_in_other_columns) — chỉ mỗi khoảng cách không đủ tin
    cậy để tự quyết định (xem docstring hàm đó), chỉ giữ lại bbox của dòng
    ĐẦU TIÊN trong nhóm làm neo hàng (row anchor)."""
    if not bboxes_sorted_by_y:
        return []
    merged = [bboxes_sorted_by_y[0]]
    for bbox in bboxes_sorted_by_y[1:]:
        prev = merged[-1]
        prev_height = max(prev[3] - prev[1], 1)
        gap = bbox[1] - prev[3]
        if gap <= prev_height * gap_factor and _label_continues_in_other_columns(
            lines, prev[3], bbox[1], label_col_x_max
        ):
            continue  # cùng 1 hàng đã wrap, bỏ qua dòng tiếp theo trong nhóm
        merged.append(bbox)
    return merged


def _detect_table(lines):
    """Thử nhận diện 1 bảng trong `lines` (từ _page_lines — DÒNG GỐC,
    KHÔNG qua merge_paragraph_blocks). Trả về None nếu không thấy, hoặc
    {"row_y": [...], "col_x": [...]} — tọa độ neo (điểm bắt đầu) của từng
    hàng/cột đã phát hiện.

    Dùng dòng gốc thay vì group đã gộp đoạn văn: với nội dung dạng bảng,
    merge_paragraph_blocks có thể dính nhiều Ô KHÁC HÀNG của cùng 1 cột
    hẹp thành 1 group duy nhất (ví dụ "Application Layer" gộp luôn với
    "Presentation Layer" thành 1 khối 2 dòng) — làm cột nhãn chỉ còn vài
    "hàng" giả thay vì đúng số hàng thật, khiến việc dò neo hàng sai lệch
    hẳn. Từng DÒNG riêng lẻ (trước khi gộp) mới phản ánh đúng ranh giới
    hàng thật của bảng.

    Cột NHÃN HÀNG (vd "Processor", "DRAM"...) dùng làm nguồn xác định số
    hàng và vị trí từng hàng, vì mỗi nhãn hàng luôn nằm gọn trong 1 dòng
    riêng, NGẮN (không wrap dài như ô dữ liệu/đoạn văn), nên đáng tin cậy
    hơn hẳn để suy ra ranh giới hàng thật của bảng. KHÔNG mặc định lấy dải
    cột trái nhất trên TRANG làm cột nhãn — dải cột trái nhất có thể là 1
    khối văn bản/tiêu đề KHÔNG liên quan tới bảng (vd đoạn mô tả nằm bên
    trái, bảng thật nằm bên phải) mà vẫn đủ ≥3 dòng để lọt qua kiểm tra số
    lượng. Thử LẦN LƯỢT từng dải cột từ trái sang phải, chấp nhận dải đầu
    tiên có nhãn NGẮN (độ dài trung vị dưới ngưỡng — đoạn văn wrap dài
    dòng nào cũng gần bằng bề rộng cột, dài hơn hẳn 1 nhãn thật) và cỡ chữ
    tương đối đồng đều."""
    if len(lines) < MIN_TABLE_ROWS * MIN_TABLE_COLS:
        return None

    col_bands = _cluster_1d(
        [(bbox[0], (bbox, size, text)) for bbox, text, size, _color in lines], TABLE_COL_GAP
    )
    col_bands = [b for b in col_bands if len(b) >= MIN_TABLE_ROWS]
    if len(col_bands) < MIN_TABLE_COLS:
        return None
    col_bands.sort(key=lambda band: min(pos for pos, _ in band))

    for label_col_idx in range(len(col_bands) - MIN_TABLE_COLS + 1):
        remaining_cols = col_bands[label_col_idx:]
        label_items = [item for _, item in remaining_cols[0]]

        # 1 nhãn/chú thích tự wrap nhiều dòng (vd "Stage 4: Layer of
        # interaction via applications" trong 1 hộp sơ đồ, KHÔNG phải
        # bảng) có thể tạo ra nhiều dòng NGẮN trông giống hệt nhãn hàng
        # thật — lọt qua được kiểm tra độ dài trung vị ở dưới. Dấu hiệu
        # phân biệt: 1 bảng thật có số dòng ở cột nhãn XẤP XỈ số dòng ở
        # (ít nhất) 1 cột dữ liệu khác (mỗi hàng đóng góp ~1 dòng cho mỗi
        # cột); nếu cột nhãn nhiều dòng vượt trội hẳn mọi cột khác, gần
        # như chắc chắn đó là 1-2 đoạn chú thích bị wrap vụn ra, không
        # phải nhiều hàng thật.
        other_counts = [len(band) for band in remaining_cols[1:]]
        if other_counts and len(label_items) > max(other_counts) * MAX_LABEL_TO_DATA_RATIO:
            continue

        lengths = sorted(len(text) for _bbox, _size, text in label_items)
        if lengths[len(lengths) // 2] > TABLE_LABEL_MAX_MEDIAN_LEN:
            continue

        sizes = sorted(size for _bbox, size, _text in label_items)
        median_size = sizes[len(sizes) // 2]
        row_anchors = [
            bbox for bbox, size, _text in label_items
            if abs(size - median_size) <= median_size * TABLE_ROW_LABEL_SIZE_TOLERANCE
        ]
        if len(row_anchors) < MIN_TABLE_ROWS:
            continue
        row_anchors.sort(key=lambda bbox: bbox[1])
        # remaining_cols luôn có >= MIN_TABLE_COLS (2) dải, đảm bảo bởi
        # phạm vi vòng lặp label_col_idx ở trên — luôn có remaining_cols[1].
        label_col_x_max = min(pos for pos, _ in remaining_cols[1])
        row_anchors = _merge_wrapped_row_labels(row_anchors, lines, label_col_x_max)
        row_anchors = _trim_row_anchor_outliers(row_anchors)
        if len(row_anchors) < MIN_TABLE_ROWS:
            continue

        # ponytail: từng có thêm 1 điều kiện chặn "khoảng cách hàng-hàng
        # không được dao động quá 2 lần" — bỏ đi vì bảng thật vẫn thường
        # có hàng cao thấp khác nhau (ô nhiều dòng hơn ô ít dòng), chặn
        # kiểu đó từ chối nhầm cả bảng hợp lệ. _trim_row_anchor_outliers ở
        # trên đã lo phần quan trọng hơn: cắt các ứng viên lạc lõng ở
        # ĐẦU/CUỐI dãy so với khoảng cách phổ biến nhất.
        row_ys = [bbox[1] for bbox in row_anchors]
        col_x = [min(pos for pos, _ in band) for band in remaining_cols]

        # Kiểm tra cuối, MẠNH nhất: bảng dữ liệu thật hầu như luôn được
        # điền ĐẦY ĐỦ (mọi ô hàng×cột đều có nội dung) — đã đo thực tế: cả
        # 3 bảng thật (ARTIK, 2 bảng OSI) đều lấp đầy 100%, trong khi mọi
        # sơ đồ/infographic bị nhận nhầm trước đó (biểu đồ tròn, cây phân
        # loại giao thức, sơ đồ luồng MQTT...) chỉ lấp đầy 50-75% do vị
        # trí các nhãn/icon rải rác không thực sự khớp thành lưới đều —
        # đây là dấu hiệu phân biệt rõ ràng nhất, tính trên cùng phép gán
        # ô sẽ dùng thật khi vẽ bảng (_build_table_cells), không phải ước
        # lượng riêng.
        cells = _build_table_cells(lines, row_ys, col_x)
        fill_ratio = len(cells) / (len(row_ys) * len(col_x))
        if fill_ratio < TABLE_MIN_FILL_RATIO:
            continue

        # Bảng càng DÀY ĐẶC (mỗi ô nhiều dòng, kiểu đoạn văn/gạch đầu dòng
        # dài) càng dễ khiến việc dò neo hàng theo vị trí sai — đã kiểm
        # chứng thực tế: 1 dòng LẠC (tiêu đề nhóm/chú thích, không phải
        # nhãn hàng thật) tình cờ rơi đúng cột nhãn vẫn có thể lọt qua mọi
        # kiểm tra ở trên (kể cả fill_ratio — vẫn lấp đầy vì nội dung dồi
        # dào), kéo lệch toàn bộ hàng phía sau. Không tìm được tín hiệu
        # nào tách được ca này khỏi bảng thật (bố cục tài liệu quá tự do,
        # nhãn hàng thật cũng có khi lệch xa nội dung cột bên) — AN TOÀN
        # HƠN là từ chối dựng bảng khi độ dày vượt ngưỡng đã kiểm chứng an
        # toàn (bảng thật đơn giản, đã test kỹ, không ô nào quá 3 dòng),
        # để nội dung vẫn được vẽ ĐÚNG VỊ TRÍ THẬT qua text box thường
        # (convert_to_pptx() tự làm việc này khi _detect_table() trả None)
        # thay vì có nguy cơ lệch hàng trong 1 bảng "trông như đúng".
        max_cell_lines = max((len(items) for items in cells.values()), default=0)
        if max_cell_lines > TABLE_MAX_CELL_LINES:
            continue

        # 2 sơ đồ giống hệt nhau đặt cạnh nhau (vd cùng 1 chồng OSI 7 lớp
        # lặp lại cho 2 thiết bị khác nhau) có thể xếp thành lưới ĐẦY ĐỦ y
        # hệt 1 bảng thật, nhưng cột nhãn và cột kế bên chỉ là 2 BẢN SAO
        # giống hệt nhau — 1 bảng dữ liệu thật luôn có nội dung KHÁC NHAU
        # giữa các cột (đó là lý do có nhiều cột). Loại nếu cột đầu và cột
        # kế trùng gần như 100% theo từng hàng.
        if len(col_x) > 1:
            col0_texts = [" ".join(t for _b, t, _s, _c in cells.get((r, 0), [])) for r in range(len(row_ys))]
            col1_texts = [" ".join(t for _b, t, _s, _c in cells.get((r, 1), [])) for r in range(len(row_ys))]
            matches = sum(1 for a, b in zip(col0_texts, col1_texts) if a and a == b)
            if matches >= len(row_ys) * 0.8:
                continue

        return {"row_y": row_ys, "col_x": col_x}

    return None


def _trim_row_anchor_outliers(bboxes_sorted_by_y, outlier_factor=2.2):
    """Cỡ chữ đôi khi không đủ khác biệt để loại 1 dòng mô tả/tiêu đề phụ
    nằm ngay TRÊN bảng ra khỏi cột nhãn (ví dụ dòng "• A comparison of..."
    cỡ chữ khá gần với nhãn hàng thật) — dòng đó vẫn tạo ra 1 khoảng cách
    NHẢY VỌT bất thường so với khoảng cách đều đặn giữa các hàng thật.
    Chỉ giữ lại đoạn LIÊN TỤC có khoảng cách hàng-hàng gần với khoảng cách
    phổ biến nhất, cắt bỏ các ứng viên lạc lõng ở đầu/cuối."""
    if len(bboxes_sorted_by_y) < MIN_TABLE_ROWS + 1:
        return bboxes_sorted_by_y
    ys = [bbox[1] for bbox in bboxes_sorted_by_y]
    gaps = [ys[i + 1] - ys[i] for i in range(len(ys) - 1)]
    median_gap = sorted(gaps)[len(gaps) // 2]
    start = 0
    while start < len(gaps) and gaps[start] > median_gap * outlier_factor:
        start += 1
    end = len(bboxes_sorted_by_y)
    while end - 2 >= 0 and gaps[end - 2] > median_gap * outlier_factor:
        end -= 1
    return bboxes_sorted_by_y[start:end]


def _line_color(line, ocr_rects):
    """Màu chữ của 1 dòng. Trang OCR luôn có color=(0,0,0) hardcode ở
    span (chữ ẩn chèn lúc OCR không biết màu thật) — phải khớp với rect
    gốc Vision nhận diện được (đã dò màu mực thật lúc OCR, xem ocr_pdf.py)
    thay vì đọc thẳng span["color"], nếu không mọi ô bảng sẽ luôn ra màu
    đen bất kể màu thật trong ảnh."""
    if ocr_rects:
        rect = fitz.Rect(line["bbox"])
        for orig_rect, ink_color in ocr_rects:
            if orig_rect.intersects(rect):
                return ink_color
    color = line["spans"][0].get("color")
    if color is None:
        return (0, 0, 0)
    return (color >> 16) & 255, (color >> 8) & 255, color & 255


def _page_lines(page, ocr_rects=None):
    """Toàn bộ dòng chữ của trang, đã tách các dòng KHÔNG liên quan mà
    PyMuPDF lỡ gộp nhầm (xem split_incoherent_block) — nhưng KHÔNG gộp
    đoạn văn (merge_paragraph_blocks): với bảng biểu, gộp đoạn văn có thể
    dính 2 giá trị thuộc 2 HÀNG khác nhau trong cùng 1 cột hẹp (ví dụ "32
    KB" của hàng DRAM và "256 KB" của hàng FLASH Memory bị nối thành 1
    dòng) — cần giữ từng dòng gốc riêng biệt để gán đúng hàng."""
    raw_blocks = page.get_text("dict")["blocks"]
    blocks = [
        sub for b in raw_blocks
        for sub in (split_incoherent_block(b) if b.get("type") == 0 else [b])
    ]
    lines = []
    for b in blocks:
        if b.get("type") != 0:
            continue
        for line in b["lines"]:
            # Bỏ qua nhãn xoay dọc (vd "Applications"/"Communications"
            # cạnh sơ đồ OSI) — ô bảng luôn là chữ NGANG, 1 nhãn xoay dọc
            # tình cờ nằm gần 1 bảng thật (cùng hàng theo Y) có thể bị gán
            # nhầm vào cột nhãn của bảng nếu không loại từ đầu; nhãn xoay
            # vẫn được vẽ đúng, đúng hướng qua đường dẫn text box thường
            # (dùng `groups`, không dùng `lines` này).
            dx, _dy = line.get("dir", (1.0, 0.0))
            if abs(dx) < 0.5:
                continue
            text = "".join(s["text"] for s in line["spans"]).strip()
            if text:
                size = line["spans"][0].get("size", 11)
                lines.append((line["bbox"], text, size, _line_color(line, ocr_rects)))
    return lines


def _interval_index(value, anchors, tolerance=0.0):
    """Chỉ số của neo LỚN NHẤT mà `value` không thấp hơn quá `tolerance` —
    tức đúng khoảng [anchors[i], anchors[i+1]) mà `value` rơi vào — KHÔNG
    phải neo GẦN NHẤT theo khoảng cách tuyệt đối (_nearest_index). Dùng
    cho cả gán HÀNG lẫn gán CỘT trong _build_table_cells(), vì cả 2 đều
    gặp chung 1 vấn đề: 1 hàng/cột có thể RỘNG/CAO hơn hẳn hàng/cột lân
    cận (ô dữ liệu wrap nhiều dòng hơn nhãn hàng ngắn của chính nó; hoặc
    tiêu đề canh GIỮA nằm xa hẳn mép trái của cột rộng) — khi đó nội dung
    ở rìa xa của hàng/cột đó lại GẦN neo BÊN CẠNH hơn theo khoảng cách
    tuyệt đối dù vẫn thuộc hàng/cột hiện tại; "gần nhất" gán sai, "khoảng
    [neo hiện tại, neo sau)" luôn đúng miễn giá trị đó CHƯA tới neo sau.
    `tolerance` bù 1 ca khác: dòng ĐẦU TIÊN của 1 ô có thể nằm trước 1 chút
    so với chính neo của nó (lệch canh dòng cơ sở giữa cột nhãn/cột dữ
    liệu) — không có tolerance sẽ bị tính nhầm thuộc hàng/cột TRƯỚC."""
    idx = 0
    for i, anchor in enumerate(anchors):
        if anchor - tolerance <= value:
            idx = i
        else:
            break
    return idx


def _build_table_cells(lines, row_y, col_x):
    """Gán mỗi dòng vào (hàng, cột) theo _interval_index() cho CẢ 2 trục
    (xem docstring ở đó) — không dùng neo GẦN NHẤT cho trục nào, vì cột
    cũng gặp đúng vấn đề như hàng: 1 tiêu đề canh GIỮA trong cột RỘNG (vd
    "Remarks" nằm cách xa hẳn mép trái cột của chính nó) lại gần mép cột
    KHÁC hơn theo khoảng cách tuyệt đối (thật sự gặp — tiêu đề bị loại
    nhầm, ô trống). Loại các dòng nằm hẳn ngoài phạm vi hàng/cột đầu/cuối
    của bảng — ví dụ 1 đoạn văn ở tận rìa trái trang lại tình cờ thẳng
    hàng dọc (Y) với 1 nhãn hàng thật dù cột (X) của nó cách xa cả bảng
    hàng trăm điểm."""
    row_gaps = [row_y[i + 1] - row_y[i] for i in range(len(row_y) - 1)]
    row_reject = (sum(row_gaps) / len(row_gaps)) * 0.9 if row_gaps else 20.0
    col_gaps = [col_x[i + 1] - col_x[i] for i in range(len(col_x) - 1)]
    col_reject = (sum(col_gaps) / len(col_gaps)) * 0.9 if col_gaps else 60.0

    cells = {}
    for bbox, text, size, color in lines:
        if bbox[1] < row_y[0] - row_reject or bbox[1] > row_y[-1] + row_reject * 3:
            continue
        if bbox[0] < col_x[0] - col_reject or bbox[0] > col_x[-1] + col_reject * 3:
            continue
        row_idx = _interval_index(bbox[1], row_y, TABLE_ROW_BASELINE_TOLERANCE)
        col_idx = _interval_index(bbox[0], col_x)
        cells.setdefault((row_idx, col_idx), []).append((bbox, text, size, color))
    for key, items in cells.items():
        items.sort(key=lambda item: item[0][1])
    return cells


def _table_crosses_section_divider(page, table_bbox_rect, min_width_ratio=0.9, min_height=5.0):
    """True nếu có 1 thanh/khối trang trí gần như CẢ CHIỀU RỘNG TRANG
    (>= min_width_ratio) VÀ đủ cao (>= min_height — loại các vạch kẻ mỏng
    bình thường giữa các dòng/hàng, chỉ ~0.5pt, không phải ranh giới nhóm)
    cắt ngang qua vùng DỌC của bảng — dấu hiệu bảng đang gộp nhầm nội dung
    của 2 NHÓM khác nhau (được tác giả tài liệu ngăn cách bằng 1 thanh
    tiêu đề nhóm, vd nền xám đậm "Giải pháp X", hoặc cả 1 khối nền của
    nhóm KHÁC) thành 1 bảng liên tục, vì _detect_table() chỉ dò theo 1 cột
    nhãn hẹp, không biết gì về ranh giới nhóm này — bug thật đã gặp: 1
    bảng nhỏ (chỉ 2 cột: tên sản phẩm + ghi chú) vô tình trải dài qua
    đúng 1 thanh ngăn cách như vậy, kéo theo cả những hàng thuộc nhóm
    giải pháp KHÁC vào chung 1 bảng.

    An toàn hơn là từ chối hẳn bảng này (xem TABLE_MAX_CELL_LINES) thay vì
    cố tách lại đúng ranh giới — bố cục tài liệu quá đa dạng để làm việc
    đó đáng tin cậy."""
    page_w = page.rect.width
    y0, y1 = table_bbox_rect.y0, table_bbox_rect.y1
    for d in page.get_drawings():
        r = d.get("rect")
        if (
            r
            and r.width >= page_w * min_width_ratio
            and r.height >= min_height
            and y0 + 2 < r.y0
            and r.y1 < y1 - 2
        ):
            return True
    return False


def _table_bbox_rect(table_info, cells):
    """Vùng hình chữ nhật (left, top, width, height, tính bằng point) mà
    _add_table_shape() sẽ đặt bảng PowerPoint vào — trích riêng ra khỏi
    _add_table_shape() để convert_to_pptx() gọi được TRƯỚC khi vẽ ảnh nền
    (cần biết trước vùng này để xóa vạch kẻ/nền vector THẬT của bảng gốc
    khỏi ảnh nền, xem _page_background_png()), tránh tính 2 lần 2 công
    thức có thể lệch nhau. Trả về None nếu `cells` rỗng (không có gì để
    dựng bảng)."""
    row_y, col_x = table_info["row_y"], table_info["col_x"]
    rows, cols = len(row_y), len(col_x)
    all_items = [item for group in cells.values() for item in group]
    if not all_items:
        return None
    col_step = (col_x[-1] - col_x[0]) / max(cols - 1, 1) if cols > 1 else 120.0
    row_step = (row_y[-1] - row_y[0]) / max(rows - 1, 1) if rows > 1 else 20.0
    left, top = col_x[0] - 4, row_y[0] - 4
    # Biên phải/dưới của bảng: dùng bbox THẬT xa nhất của các dòng đã gán
    # vào bảng (không chỉ nới thêm 1 khoảng cách cột/hàng trung bình vào
    # điểm neo cuối) — nền trang scan có thể có 1 dải màu trang trí (như
    # sơ đồ OSI) rộng hơn hẳn ước lượng ngoại suy đơn giản; nếu bảng vẽ
    # đè lên hẹp hơn dải màu gốc, phần dải màu còn lại sẽ lộ ra ngoài rìa
    # bảng, trông như 1 mảng màu lạc quẻ không thuộc bảng.
    max_x1 = max(item[0][2] for item in all_items)
    max_y1 = max(item[0][3] for item in all_items)
    total_w = max(col_x[-1] + col_step - left, max_x1 + 4 - left)
    total_h = max(row_y[-1] + row_step - top, max_y1 + 4 - top)
    return left, top, total_w, total_h


def _add_table_shape(slide, table_info, cells):
    row_y, col_x = table_info["row_y"], table_info["col_x"]
    rows, cols = len(row_y), len(col_x)

    bbox = _table_bbox_rect(table_info, cells)
    if bbox is None:
        return
    left, top, total_w, total_h = bbox

    shape = slide.shapes.add_table(
        rows, cols,
        Emu(int(left * EMU_PER_POINT)), Emu(int(top * EMU_PER_POINT)),
        Emu(int(total_w * EMU_PER_POINT)), Emu(int(total_h * EMU_PER_POINT)),
    )
    table = shape.table

    # col_edges/row_edges PHẢI bắt đầu từ `left`/`top` (biên thật của bảng),
    # KHÔNG phải từ col_x[0]/row_y[0] — bug thật đã gặp: left = col_x[0]-4
    # (đệm 4pt sang trái/trên so với neo cột/hàng đầu tiên), nếu col_edges
    # bắt đầu từ col_x[0] thì 4pt đệm đó KHÔNG được cộng vào độ rộng cột 0
    # nào cả — tổng độ rộng/cao thật của bảng (python-pptx tự tính lại
    # bằng tổng độ rộng từng cột/hàng khi set riêng từng cái) hụt mất đúng
    # 4pt so với total_w/total_h đã dùng để xóa nền — để lại 1 viền hở lộ
    # nền/vạch kẻ gốc chưa xóa, nhìn như "bảng trong bảng".
    col_edges = [left] + list(col_x[1:]) + [left + total_w]
    row_edges = [top] + list(row_y[1:]) + [top + total_h]
    for c in range(cols):
        table.columns[c].width = Emu(int((col_edges[c + 1] - col_edges[c]) * EMU_PER_POINT))
    for r in range(rows):
        table.rows[r].height = Emu(int((row_edges[r + 1] - row_edges[r]) * EMU_PER_POINT))

    for (r, c), items in cells.items():
        cell = table.cell(r, c)
        cell.text = " ".join(text for _bbox, text, _size, _color in items)
        cell.margin_left = cell.margin_right = Emu(18000)
        cell.margin_top = cell.margin_bottom = Emu(9000)
        font_size = Pt(max(min(size for _bbox, _text, size, _color in items), 1))
        r_, g_, b_ = items[0][3]
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = font_size
                run.font.color.rgb = RGBColor(r_, g_, b_)


def convert_to_pptx(input_pdf, output_pptx):
    doc = fitz.open(input_pdf)          # nguồn để trích khối chữ
    # Giữ rect gốc Vision nhận diện (không chỉ biết trang nào được OCR) để
    # _page_background_png() xóa đúng pixel chữ cũ trên bg_doc bên dưới.
    ocr_rects_by_page = ensure_text_layer(doc)
    bg_doc = fitz.open(input_pdf)       # bản dùng-rồi-bỏ: sẽ bị xóa hết chữ (ảnh/vector giữ nguyên)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # layout "Blank"
    # 1 presentation chỉ có 1 kích thước slide chung cho tất cả (khác PDF,
    # vốn mỗi trang có thể khác kích thước); lấy kích thước trang đầu tiên
    # làm kích thước chung cho cả bộ slide.
    # ponytail: PDF có nhiều trang khác kích thước sẽ bị các trang sau co
    # theo kích thước này thay vì giữ đúng kích thước riêng — hiếm gặp
    # trong thực tế, và python-pptx cũng không hỗ trợ kích thước riêng
    # từng slide nên không đáng để xử lý thêm.
    first_page = doc[0]
    prs.slide_width = Emu(int(first_page.rect.width * EMU_PER_POINT))
    prs.slide_height = Emu(int(first_page.rect.height * EMU_PER_POINT))

    for page_index, page in enumerate(doc):
        page_w_pt, page_h_pt = page.rect.width, page.rect.height
        slide = prs.slides.add_slide(blank_layout)

        # Mọi ảnh nhúng rời rạc (kể cả ảnh lớn — xem
        # _page_image_rects_for_export()) được vẽ RIÊNG thành picture shape
        # sửa được, thay vì bị khoá cứng vào ảnh nền toàn trang như trước.
        image_rects = _page_image_rects_for_export(page, page_index in ocr_rects_by_page)

        # merge_paragraph_blocks tự gộp các dòng wrap của cùng 1 đoạn/bullet
        # thành 1 nhóm (và tự tách trước các dòng KHÔNG liên quan mà
        # PyMuPDF lỡ gộp nhầm — xem translator_engine.py) — dùng lại y hệt
        # logic đã kiểm chứng ở pipeline dịch, để 1 đoạn văn wrap nhiều dòng
        # ra đúng 1 text box, cùng 1 cỡ chữ, thay vì mỗi dòng 1 box riêng.
        raw_blocks = page.get_text("dict")["blocks"]
        groups = merge_paragraph_blocks(raw_blocks)
        lines = _page_lines(page, ocr_rects_by_page.get(page_index))

        # Phải biết TRƯỚC có dựng bảng hay không (và dựng ở đâu) rồi mới
        # vẽ ảnh nền — _page_background_png() cần table_bbox_rect để xóa
        # luôn vạch kẻ/nền vector THẬT của bảng gốc trong đúng vùng đó
        # (xem docstring ở đó), nếu không bảng PowerPoint dựng sau sẽ đè
        # lên trên 1 bảng gốc vẫn còn nguyên bên dưới.
        table_info = _detect_table(lines)
        table_bbox = None
        table_cells = None
        table_bbox_rect = None
        if table_info:
            row_y, col_x = table_info["row_y"], table_info["col_x"]
            col_step = (col_x[-1] - col_x[0]) / max(len(col_x) - 1, 1) if len(col_x) > 1 else 120.0
            row_step = (row_y[-1] - row_y[0]) / max(len(row_y) - 1, 1) if len(row_y) > 1 else 20.0
            table_bbox = (
                col_x[0] - col_step * 0.5, row_y[0] - row_step * 0.5,
                col_x[-1] + col_step * 1.5, row_y[-1] + row_step * 1.5,
            )
            table_cells = _build_table_cells(lines, row_y, col_x)
            precise_bbox = _table_bbox_rect(table_info, table_cells) if table_cells else None
            if precise_bbox:
                left, top, width, height = precise_bbox
                table_bbox_rect = fitz.Rect(left, top, left + width, top + height)
                if _table_crosses_section_divider(page, table_bbox_rect):
                    table_bbox = None
                    table_cells = None
                    table_bbox_rect = None
            else:
                table_bbox = None
                table_cells = None

        png_bytes = _page_background_png(
            bg_doc, page_index, ocr_rects_by_page, image_rects, table_bbox=table_bbox_rect
        )
        slide.shapes.add_picture(
            io.BytesIO(png_bytes), Emu(0), Emu(0),
            width=Emu(int(page_w_pt * EMU_PER_POINT)),
            height=Emu(int(page_h_pt * EMU_PER_POINT)),
        )
        for rect in image_rects:
            crop_bytes = _crop_image_png(page, rect)
            slide.shapes.add_picture(
                io.BytesIO(crop_bytes),
                Emu(int(rect.x0 * EMU_PER_POINT)), Emu(int(rect.y0 * EMU_PER_POINT)),
                width=Emu(int(rect.width * EMU_PER_POINT)),
                height=Emu(int(rect.height * EMU_PER_POINT)),
            )

        if table_cells:
            _add_table_shape(slide, table_info, table_cells)

        for group in groups:
            text = "".join(
                s["text"] for b in group["sub_blocks"] for l in b["lines"] for s in l["spans"]
            ).strip()
            if not text:
                continue
            if table_bbox is not None:
                gx0, gy0, gx1, gy1 = group["bbox"]
                cx, cy = (gx0 + gx1) / 2, (gy0 + gy1) / 2
                if table_bbox[0] <= cx <= table_bbox[2] and table_bbox[1] <= cy <= table_bbox[3]:
                    continue  # đã hiển thị trong bảng, không vẽ text box rời nữa
            _add_text_box(slide, group, page_h_pt, ocr_rects_by_page.get(page_index))

    prs.save(output_pptx)
    doc.close()
    bg_doc.close()
