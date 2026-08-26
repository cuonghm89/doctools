"""Tự kiểm tra dịch trực tiếp .docx/.pptx (giữ nguyên định dạng). Chạy trực
tiếp: python3 test_office_translate.py
Không gọi mạng: router._deepl_batch được stub (giả lập) ra, giống hệt cách
test_router.py/test_pdf_convert.py đã làm.
"""
import os
import tempfile

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE

from office_translate import (
    translate_docx, translate_pptx, FONT_SHRINK_MIN_SCALE, FONT_SHRINK_MIN_SCALE_TABLE,
    _wrapped_line_count,
)
from router import TranslationRouter


def _all_texts(shapes):
    """Đệ quy vào group shape, giống hệt _iter_pptx_text_frames trong
    office_translate.py — để test xác nhận đúng những gì hàm dịch thật đã
    xử lý, kể cả text nằm trong group (add_group_shape di chuyển shape gốc
    vào bên trong group, không còn ở top-level slide.shapes nữa)."""
    texts = set()
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            texts |= _all_texts(shape.shapes)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame.text:
                        texts.add(cell.text_frame.text)
        elif shape.has_text_frame and shape.text_frame.text:
            texts.add(shape.text_frame.text)
    return texts


def make_stub_router():
    router = TranslationRouter(deepl_key="fakekey:fx", gemini_key="")
    router.tracker = {"chars_used": 0, "limit": 500_000}
    router._deepl_batch = lambda texts: [t.upper() for t in texts]
    return router


# --- .docx: đoạn văn 2 run (in đậm + thường), 1 bảng 1 ô, header. ---
with tempfile.TemporaryDirectory() as tmp:
    docx_in = os.path.join(tmp, "in.docx")
    docx_out = os.path.join(tmp, "out.docx")

    doc = Document()
    doc.sections[0].header.paragraphs[0].add_run("Header line here")
    p = doc.add_paragraph()
    run1 = p.add_run("Bold start ")
    run1.bold = True
    run1.font.size = Pt(14)
    run2 = p.add_run("normal end.")
    table = doc.add_table(rows=1, cols=1)
    table.rows[0].cells[0].paragraphs[0].add_run("Cell content")
    doc.save(docx_in)

    # Stub trả về NGUYÊN VĂN bản gốc (bề rộng bằng 0% chênh lệch) cho test
    # này — mục đích ở đây là kiểm tra ĐỊNH DẠNG (in đậm, cỡ chữ) được giữ
    # nguyên khi bề rộng bản dịch KHÔNG phình, tách riêng khỏi test co-chữ-
    # khi-phình-dài ở dưới (dùng router riêng, cố tình phình rất dài).
    format_router = TranslationRouter(deepl_key="fakekey:fx", gemini_key="")
    format_router.tracker = {"chars_used": 0, "limit": 500_000}
    format_router._deepl_batch = lambda texts: list(texts)
    translate_docx(docx_in, docx_out, format_router)

    out_doc = Document(docx_out)
    body_para = out_doc.paragraphs[0]  # đoạn văn add_paragraph() đầu tiên trong body
    assert body_para.text == "Bold start normal end.", body_para.text
    assert body_para.runs[0].bold is True, "định dạng run đầu (in đậm) phải giữ nguyên"
    assert body_para.runs[0].font.size == Pt(14), "cỡ chữ run đầu phải giữ nguyên khi bề rộng bản dịch không phình"
    assert all(r.text == "" for r in body_para.runs[1:]), "các run sau phải bị xóa rỗng"

    assert out_doc.tables[0].rows[0].cells[0].text == "Cell content"
    assert out_doc.sections[0].header.paragraphs[0].text == "Header line here"

print("Tất cả self-check dịch .docx đều pass.")


# --- .pptx: text box thường, bảng, group shape, và max_pages chỉ dịch N
# slide đầu (slide sau max_pages phải giữ NGUYÊN bản gốc). ---
with tempfile.TemporaryDirectory() as tmp2:
    pptx_in = os.path.join(tmp2, "in.pptx")
    pptx_out = os.path.join(tmp2, "out.pptx")

    prs = Presentation()
    blank = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank)
    tb = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tb.text_frame.paragraphs[0].add_run().text = "Slide one text"
    table_shape = slide1.shapes.add_table(1, 1, Inches(1), Inches(3), Inches(2), Inches(1))
    table_shape.table.rows[0].cells[0].text = "Table cell"
    group = slide1.shapes.add_group_shape([tb])  # group chứa lại chính tb -> vẫn phải dịch được qua đệ quy

    slide2 = prs.slides.add_slide(blank)
    tb2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tb2.text_frame.paragraphs[0].add_run().text = "Slide two text"

    prs.save(pptx_in)

    translate_pptx(pptx_in, pptx_out, make_stub_router(), max_pages=1)

    out_prs = Presentation(pptx_out)
    slide1_texts = _all_texts(out_prs.slides[0].shapes)
    assert "SLIDE ONE TEXT" in slide1_texts, slide1_texts
    assert "TABLE CELL" in slide1_texts, "phải dịch được cả ô bảng"

    slide2_texts = _all_texts(out_prs.slides[1].shapes)
    assert "Slide two text" in slide2_texts, f"quá max_pages=1 phải giữ nguyên bản gốc, có {slide2_texts}"

print("Tất cả self-check dịch .pptx đều pass.")


# --- Bản dịch phình dài hơn hẳn bản gốc phải làm cỡ chữ GIẢM xuống (chặn
# dưới FONT_SHRINK_MIN_SCALE) — tránh tràn/chồng lấn lên shape khác. Text
# box cũng phải bật auto_size (TEXT_TO_FIT_SHAPE) để PowerPoint thật tự co
# thêm khi cần. Router riêng cho test này dịch ra chuỗi dài gấp 3, để chắc
# chắn vượt ngưỡng kích hoạt co chữ. ---
def make_long_stub_router():
    router = TranslationRouter(deepl_key="fakekey:fx", gemini_key="")
    router.tracker = {"chars_used": 0, "limit": 500_000}
    router._deepl_batch = lambda texts: [t + " " + t + " " + t for t in texts]
    return router


with tempfile.TemporaryDirectory() as tmp3:
    pptx_in3 = os.path.join(tmp3, "in.pptx")
    pptx_out3 = os.path.join(tmp3, "out.pptx")

    prs3 = Presentation()
    slide3 = prs3.slides.add_slide(prs3.slide_layouts[6])
    tb3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    run3 = tb3.text_frame.paragraphs[0].add_run()
    run3.text = "Short title"
    run3.font.size = PptxPt(24)
    prs3.save(pptx_in3)

    translate_pptx(pptx_in3, pptx_out3, make_long_stub_router())

    out_prs3 = Presentation(pptx_out3)
    out_tb = out_prs3.slides[0].shapes[0]
    assert out_tb.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, \
        "text box phải bật auto_size để PowerPoint tự co thêm khi mở file"
    out_run = out_tb.text_frame.paragraphs[0].runs[0]
    assert out_run.text == "Short title Short title Short title", (
        f"cụm ngắn (<=4 từ) KHÔNG được bỏ dịch dù bản dịch dài gấp 3 — "
        f"guard đó chỉ đúng cho PDF (khung cố định), ra {out_run.text!r}"
    )
    assert out_run.font.size < PptxPt(24), f"cỡ chữ phải giảm khi bản dịch dài hơn hẳn, ra {out_run.font.size}"
    assert out_run.font.size >= PptxPt(24 * 0.55), "không được co quá FONT_SHRINK_MIN_SCALE"

print("Tất cả self-check co-chữ-khi-phình-dài đều pass.")


# --- Ô bảng chật hơn text box tự do -> ngưỡng co tối thiểu RIÊNG, THẤP
# hơn (FONT_SHRINK_MIN_SCALE_TABLE=0.40 so với 0.55). word_wrap phải bị ép
# bật kể cả khi file gốc tắt (tránh tràn ngang sang cột bên). Cùng 1 mức
# phình (~2.16x đo bằng font thật) áp cho cả text box lẫn ô bảng để so
# sánh trực tiếp: ô bảng phải co NHIỀU HƠN (không bị chặn sớm ở 0.55). ---
def make_double_stub_router():
    router = TranslationRouter(deepl_key="fakekey:fx", gemini_key="")
    router.tracker = {"chars_used": 0, "limit": 500_000}
    router._deepl_batch = lambda texts: [t + " " + t for t in texts]
    return router


with tempfile.TemporaryDirectory() as tmp4:
    pptx_in4 = os.path.join(tmp4, "in.pptx")
    pptx_out4 = os.path.join(tmp4, "out.pptx")

    prs4 = Presentation()
    slide4 = prs4.slides.add_slide(prs4.slide_layouts[6])
    tb4 = slide4.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    run4 = tb4.text_frame.paragraphs[0].add_run()
    run4.text = "Cell"
    run4.font.size = PptxPt(24)

    table4 = slide4.shapes.add_table(1, 1, Inches(1), Inches(3), Inches(2), Inches(1)).table
    table4.cell(0, 0).text_frame.word_wrap = False  # cố tình tắt -> phải bị ép bật lại
    cell_run = table4.cell(0, 0).text_frame.paragraphs[0].add_run()
    cell_run.text = "Cell"
    cell_run.font.size = PptxPt(24)

    prs4.save(pptx_in4)
    translate_pptx(pptx_in4, pptx_out4, make_double_stub_router())

    out_prs4 = Presentation(pptx_out4)
    out_shapes = out_prs4.slides[0].shapes
    out_tb4 = next(s for s in out_shapes if s.has_text_frame)
    out_table4 = next(s for s in out_shapes if s.has_table).table

    assert out_table4.cell(0, 0).text_frame.word_wrap is True, "word_wrap phải bị ép bật cho ô bảng"

    textbox_size = out_tb4.text_frame.paragraphs[0].runs[0].font.size
    cell_size = out_table4.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size
    assert textbox_size == PptxPt(24 * FONT_SHRINK_MIN_SCALE), (
        f"text box phải bị chặn đúng ở ngưỡng 0.55, ra {textbox_size}"
    )
    assert cell_size < textbox_size, (
        f"ô bảng phải co NHIỀU HƠN text box (ngưỡng 0.40 thấp hơn 0.55) với cùng mức phình, "
        f"cell={cell_size} textbox={textbox_size}"
    )

print("Tất cả self-check ngưỡng-co-riêng-cho-bảng đều pass.")


# --- Bản dịch phình QUÁ nhiều (co chữ tới ngưỡng 0.40 vẫn không đủ chỗ)
# -> hàng phải GIÃN CAO hơn để chứa hết, không còn tràn/chồng lấn hàng
# dưới nữa — lớp phòng thủ thứ 2, đảm bảo chắc hơn hẳn chỉ co chữ. ---
LONG_TRANSLATION = (
    "This is a much much much longer translated sentence than the original short label"
)


def make_extreme_stub_router():
    router = TranslationRouter(deepl_key="fakekey:fx", gemini_key="")
    router.tracker = {"chars_used": 0, "limit": 500_000}
    router._deepl_batch = lambda texts: [LONG_TRANSLATION for _ in texts]
    return router


with tempfile.TemporaryDirectory() as tmp5:
    pptx_in5 = os.path.join(tmp5, "in.pptx")
    pptx_out5 = os.path.join(tmp5, "out.pptx")

    prs5 = Presentation()
    slide5 = prs5.slides.add_slide(prs5.slide_layouts[6])
    table5 = slide5.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(1), Inches(0.3)).table
    original_row_height = table5.rows[0].height
    cell_run5 = table5.cell(0, 0).text_frame.paragraphs[0].add_run()
    cell_run5.text = "Data"  # không phải "A" (khớp regex acronym -> bị skip, không gọi dịch)
    cell_run5.font.size = PptxPt(24)
    prs5.save(pptx_in5)

    translate_pptx(pptx_in5, pptx_out5, make_extreme_stub_router())

    out_prs5 = Presentation(pptx_out5)
    out_table5 = next(s for s in out_prs5.slides[0].shapes if s.has_table).table
    out_cell5 = out_table5.cell(0, 0)
    out_run5 = out_cell5.text_frame.paragraphs[0].runs[0]

    assert out_run5.font.size == PptxPt(24 * FONT_SHRINK_MIN_SCALE_TABLE), (
        f"phình quá nhiều phải bị co chạm ngưỡng tối thiểu của bảng, ra {out_run5.font.size}"
    )

    usable_width_pt = (out_table5.columns[0].width - 2 * 91440) / 12700
    expected_lines = _wrapped_line_count(LONG_TRANSLATION, out_run5.font.size.pt, usable_width_pt)
    assert expected_lines > 1, "case test phải thật sự cần nhiều hơn 1 dòng mới có ý nghĩa"
    assert out_table5.rows[0].height > original_row_height, (
        f"hàng phải được GIÃN CAO hơn khi co chữ vẫn không đủ chỗ (cần {expected_lines} dòng), "
        f"gốc={original_row_height} mới={out_table5.rows[0].height}"
    )

print("Tất cả self-check giãn-hàng-khi-co-chữ-không-đủ đều pass.")
